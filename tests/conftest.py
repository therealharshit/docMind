from pathlib import Path

import pytest
from pptx import Presentation


@pytest.fixture()
def temp_env(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> Path:
    storage = tmp_path / "storage"
    monkeypatch.setenv("STORAGE_DIR", str(storage))
    monkeypatch.setenv("DATABASE_URL", f"sqlite:///{storage / 'app.db'}")
    monkeypatch.setenv("OLLAMA_BASE_URL", "http://localhost:9")
    from app.core.config import get_settings

    get_settings.cache_clear()
    yield storage
    get_settings.cache_clear()


@pytest.fixture()
def sample_pdf(tmp_path: Path) -> Path:
    import fitz

    path = tmp_path / "sample.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text(
        (72, 72),
        "Quarterly Plan\nRevenue grew 20 percent.\nCustomer churn decreased.",
    )
    doc.save(path)
    doc.close()
    return path


@pytest.fixture()
def sample_pptx(tmp_path: Path) -> Path:
    path = tmp_path / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Roadmap"
    slide.placeholders[1].text = "Launch ingestion API\nAdd parser benchmarks"
    notes = slide.notes_slide.notes_text_frame
    notes.text = "Explain the ingestion workflow and benchmark target."
    prs.save(path)
    return path
