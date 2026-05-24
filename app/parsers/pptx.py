from pathlib import Path
from xml.etree import ElementTree
from zipfile import BadZipFile, ZipFile

from pptx import Presentation

from app.core.errors import ErrorCode, ParserError
from app.parsers.base import DocumentParser
from app.schemas.document import (
    DocumentMetadata,
    DocumentType,
    ExtractedImage,
    FinalDocument,
    ProcessingMode,
    Provenance,
    Slide,
)

TEXT_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


class PPTXParser(DocumentParser):
    """Extract slide text, image metadata, and speaker notes from PPTX files."""

    def parse(
        self,
        document_id: str,
        filename: str,
        path: Path,
        processing_mode: ProcessingMode,
    ) -> FinalDocument:
        try:
            presentation = Presentation(path)
        except Exception as exc:
            raise ParserError(ErrorCode.CORRUPT_DOCUMENT, "PPTX could not be opened.") from exc

        notes_by_slide, notes_warnings = self._extract_notes(path)
        slides: list[Slide] = []
        images: list[ExtractedImage] = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            title: str | None = None
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text.strip()
                    if text:
                        if title is None:
                            title = text.splitlines()[0].strip() or None
                        texts.append(text)
                if getattr(shape, "shape_type", None) == 13:
                    image = shape.image
                    images.append(
                        ExtractedImage(
                            image_id=f"{document_id}-s{slide_index}-i{len(images) + 1}",
                            source_type="pptx",
                            slide_number=slide_index,
                            extension=image.ext,
                            width=getattr(shape, "width", None),
                            height=getattr(shape, "height", None),
                        )
                    )

            slides.append(
                Slide(
                    slide_number=slide_index,
                    title=title,
                    body_text="\n".join(texts),
                    notes=notes_by_slide.get(slide_index),
                    provenance=[
                        Provenance(
                            source_type="pptx_slide",
                            slide_number=slide_index,
                            confidence=1.0,
                        )
                    ],
                    diagnostics={"notes_extracted": slide_index in notes_by_slide},
                )
            )

        return FinalDocument(
            document_metadata=DocumentMetadata(
                document_id=document_id,
                filename=filename,
                document_type=DocumentType.PPTX,
                slide_count=len(slides),
                parser="python-pptx+xml-notes",
                processing_mode=processing_mode,
                diagnostics={
                    "ocr_skipped": True,
                    "ocr_status": "deferred_in_first_slice",
                    "notes_warnings": notes_warnings,
                },
            ),
            slides=slides,
            images=images,
        )

    def _extract_notes(self, path: Path) -> tuple[dict[int, str], list[str]]:
        try:
            with ZipFile(path) as archive:
                note_names = sorted(_note_names(archive))
                notes: dict[int, str] = {}
                warnings: list[str] = []
                for note_name in note_names:
                    slide_number = self._slide_number_from_note_name(note_name)
                    if slide_number is None:
                        warnings.append(f"Could not map note file {note_name} to slide number.")
                        continue
                    try:
                        root = ElementTree.fromstring(archive.read(note_name))
                    except ElementTree.ParseError:
                        warnings.append(f"Malformed notes XML in {note_name}.")
                        continue
                    text_parts = [
                        node.text.strip()
                        for node in root.findall(".//a:t", TEXT_NS)
                        if node.text and node.text.strip()
                    ]
                    if text_parts:
                        notes[slide_number] = "\n".join(text_parts)
                return notes, warnings
        except BadZipFile as exc:
            raise ParserError(ErrorCode.CORRUPT_DOCUMENT, "PPTX package is corrupt.") from exc

    def _slide_number_from_note_name(self, note_name: str) -> int | None:
        stem = Path(note_name).stem
        digits = "".join(char for char in stem if char.isdigit())
        return int(digits) if digits else None


def _note_names(archive: ZipFile) -> list[str]:
    return [
        name
        for name in archive.namelist()
        if name.startswith("ppt/notesSlides/notesSlide")
    ]
